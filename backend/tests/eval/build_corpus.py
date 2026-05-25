"""Genera el corpus de evaluación en distintos formatos.

Formatos creados (12 ficheros):
  Texto/Office moderno:
    sistema_solar.txt           (texto plano)
    python_lenguaje.md          (Markdown)
    cambio_climatico.pdf        (PDF via PyMuPDF)
    inteligencia_artificial.docx (Word moderno via python-docx)
    historia_internet.pptx      (PowerPoint moderno via python-pptx)
    energias_renovables.xlsx    (Excel via openpyxl)
  OpenDocument:
    biologia_celular.odt        (OpenDocument Text via odfpy)
    tabla_periodica.ods         (OpenDocument Spreadsheet via odfpy)
    sistema_nervioso.odp        (OpenDocument Presentation via odfpy)
  Formatos binarios / especiales:
    historia_fisica.rtf         (Rich Text Format)
    historia_computacion.epub   (EPUB e-book, ZIP manual)
    economia_basica.doc         (Word 97 binario, OLE2 manual)
    astronomia.ppt              (PowerPoint 97 binario, OLE2 + átomos)

Ejecutar:
  python tests/eval/build_corpus.py
"""

from __future__ import annotations

import struct
import zipfile
from pathlib import Path

OUT = Path(__file__).parent / "corpus"
OUT.mkdir(exist_ok=True)

# ═══════════════════════════════════════════════════════════════════════════
# Textos fuente
# ═══════════════════════════════════════════════════════════════════════════

PYTHON_MD = """\
# Python: lenguaje de programación
*Fuente: Wikipedia*

Python es un lenguaje de programación de **alto nivel**, interpretado y de propósito general,
cuya filosofía hace hincapié en la legibilidad del código.
Fue creado por **Guido van Rossum** y su primera versión pública, Python 0.9.0, se publicó en **1991**.
El nombre del lenguaje no proviene de la serpiente, sino del grupo de comedia británico **Monty Python**.

Es administrado por la **Python Software Foundation (PSF)** con licencia de código abierto.
En 2025 Python se posicionó como el **lenguaje de programación más popular del mundo**,
con una ventaja récord de más de 15 puntos porcentuales.

## Historia

La concepción de Python comenzó a finales de los años ochenta en el Centro para las Matemáticas
y la Informática (CWI) en los Países Bajos. Python fue diseñado como sucesor del lenguaje ABC.

- **Python 2.0** se publicó el 16 de octubre de 2000 e introdujo comprensión de listas.
- **Python 3.0** se publicó el 3 de diciembre de 2008 (revisión mayor no retrocompatible).
- El soporte de Python 2 finalizó el **1 de enero de 2020**.

## Características

Python es **multiparadigma**: orientado a objetos, imperativo y funcional.
Usa **indentación significativa** para delimitar bloques de código.

Filosofía recogida en el PEP 20 ("El Zen de Python"):
> Lo explícito es mejor que lo implícito. Lo simple es mejor que lo complejo. La legibilidad cuenta.

## Usos principales

| Área                      | Herramientas destacadas              |
|---------------------------|--------------------------------------|
| Desarrollo web            | Django, Flask, FastAPI               |
| Ciencia de datos / ML     | NumPy, Pandas, TensorFlow, PyTorch   |
| Inteligencia artificial   | LangChain, Hugging Face              |
| Automatización / scripting| Ansible, Fabric                      |
| Videojuegos               | Pygame                               |

## Tipos de datos integrados

`int`, `float`, `str`, `list` (mutable, ordenada), `tuple` (inmutable, ordenada),
`set` (mutable, sin orden), `dict` (pares clave-valor).
"""

SOLAR_TXT = """\
SISTEMA SOLAR
Fuente: Wikipedia

El sistema solar es el sistema planetario que liga gravitacionalmente a un conjunto de objetos
astronómicos que orbitan alrededor de una única estrella: el Sol. La estrella concentra el
99,86 % de la masa total del sistema. Se formó hace aproximadamente 4600 millones de años.

LOS OCHO PLANETAS

Planetas terrestres: Mercurio, Venus, Tierra, Marte.
Planetas gigantes gaseosos: Júpiter, Saturno, Urano, Neptuno.

EL SOL

Estrella de tipo espectral G2V, con unos 4600 millones de años de edad.
Temperatura superficial: ~5500 °C. Núcleo: ~15 millones de °C.
La luz tarda 8 minutos y 20 segundos en llegar del Sol a la Tierra.

UNIDAD ASTRONÓMICA

La distancia Tierra-Sol se denomina Unidad Astronómica (UA): ~149,6 millones de kilómetros.

OTROS CUERPOS

Cinturón de asteroides: entre Marte y Júpiter.
Plutón: reclasificado como planeta enano en 2006 por la UAI.
Cinturón de Kuiper: más allá de Neptuno, cuerpos helados.
Nube de Oort: reservorio esférico de cometas en los confines del sistema solar.
"""

CLIMA_TXT = """\
Cambio Climático
Fuente: Wikipedia

El cambio climático es la variación a largo plazo de los patrones climáticos globales.
El término se usa hoy principalmente para el calentamiento global antropogénico causado
por la actividad humana desde la Revolución Industrial.

CAUSAS

La principal causa es el incremento de gases de efecto invernadero (GEI):
- Dióxido de carbono (CO2): quema de combustibles fósiles y deforestación.
- Metano (CH4): ganadería, vertederos, extracción de gas.
- Óxido nitroso (N2O): agricultura y fertilizantes.

La concentración de CO2 ha pasado de 280 ppm (era preindustrial) a más de 420 ppm en 2023.

CONSECUENCIAS

- Temperatura: aumento de ~1,1 °C desde la era preindustrial.
- Nivel del mar: subido ~20 cm desde 1900.
- Fenómenos extremos: huracanes, sequías, inundaciones y olas de calor más frecuentes.
- Acidificación de océanos: amenaza arrecifes de coral.

EL ACUERDO DE PARIS (2015)

196 países firmaron el Acuerdo de París para mantener el aumento de temperatura
por debajo de 2 °C respecto a niveles preindustriales.
"""

IA_CONTENT = [
    ("Inteligencia Artificial", """\
La inteligencia artificial (IA) es la disciplina dentro de las ciencias de la computación
que busca crear sistemas capaces de realizar tareas que, cuando las hace un ser humano,
requieren inteligencia. El término fue acuñado en 1956 por John McCarthy en la Conferencia
de Dartmouth, considerada el nacimiento formal de la IA como campo de investigación."""),
    ("Hitos históricos", """\
1950 — Alan Turing publica «Computing Machinery and Intelligence» y propone el Test de Turing.
1956 — John McCarthy acuña «inteligencia artificial» en la Conferencia de Dartmouth.
1997 — Deep Blue de IBM derrota al campeón mundial de ajedrez Garry Kasparov.
2016 — AlphaGo de DeepMind derrota al campeón mundial de Go Lee Sedol.
2022 — ChatGPT de OpenAI alcanza 100 millones de usuarios en dos meses."""),
    ("Subcampos principales", """\
Aprendizaje automático (Machine Learning): sistemas que aprenden de datos.
Aprendizaje profundo (Deep Learning): redes neuronales multicapa.
Procesamiento del lenguaje natural (NLP): entender y generar lenguaje humano.
Visión por computadora: interpretación de imágenes y vídeo.
Sistemas expertos: emulan la toma de decisiones de un experto."""),
]

INTERNET_PPTX = [
    ("Historia de Internet", "De ARPANET a la red global"),
    ("Orígenes: ARPANET (1969)",
     "ARPANET fue financiado por el Departamento de Defensa de EE. UU. (DARPA).\n"
     "El 29 de octubre de 1969 se envió el primer mensaje entre UCLA y Stanford.\n"
     "J.C.R. Licklider del MIT describió en 1962 la 'Red Galáctica'."),
    ("La World Wide Web (1989-1991)",
     "Tim Berners-Lee, científico del CERN en Suiza, propuso la Web en 1989.\n"
     "En 1991 publicó la primera página web del mundo.\n"
     "Inventó HTML, HTTP y el sistema de URLs."),
    ("Internet hoy",
     "2023: ~5300 millones de usuarios (66 % de la población mundial).\n"
     "IPv6 permite 3,4 × 10^38 direcciones frente a los ~4300 M de IPv4."),
]

# ── Nuevos contenidos ──────────────────────────────────────────────────────

BIOLOGIA_ODT = [
    ("Biología Celular",
     "La biología celular es la rama de la biología que estudia la célula como unidad "
     "fundamental de la vida. Toda célula procede de otra célula preexistente, según el "
     "principio formulado por Rudolf Virchow en 1855."),
    ("Tipos de células",
     "Células procariotas: sin núcleo definido. Incluyen bacterias y arqueas. "
     "Son las formas de vida más antiguas, con unos 3500 millones de años de antigüedad.\n"
     "Células eucariotas: con núcleo rodeado de membrana nuclear. Forman todos los organismos "
     "multicelulares: plantas, animales, hongos y protistas."),
    ("El ADN",
     "El ácido desoxirribonucleico (ADN) es la molécula que almacena la información genética. "
     "Se organiza en cromosomas dentro del núcleo celular. El genoma humano contiene "
     "aproximadamente 3000 millones de pares de bases y unos 20 000 genes codificantes."),
    ("La mitosis",
     "La mitosis es el proceso de división celular en células eucariotas que produce dos "
     "células hijas con el mismo número de cromosomas que la célula madre. Consta de cuatro "
     "fases: profase, metafase, anafase y telofase, seguidas de la citocinesis."),
    ("Organelas principales",
     "Mitocondria: produce energía en forma de ATP mediante la respiración celular. "
     "Ribosomas: sintetizan proteínas. "
     "Retículo endoplasmático: transporte y modificación de proteínas y lípidos. "
     "Aparato de Golgi: empaqueta y distribuye proteínas. "
     "Lisosomas: digestión intracelular mediante enzimas hidrolíticas."),
]

TABLA_ODS = {
    "Tabla Periódica (fragmento)": [
        ["Número atómico", "Nombre", "Símbolo", "Masa atómica (u)", "Grupo", "Período"],
        ["1",  "Hidrógeno",  "H",  "1,008",  "1",  "1"],
        ["2",  "Helio",      "He", "4,003",  "18", "1"],
        ["3",  "Litio",      "Li", "6,941",  "1",  "2"],
        ["6",  "Carbono",    "C",  "12,011", "14", "2"],
        ["7",  "Nitrógeno",  "N",  "14,007", "15", "2"],
        ["8",  "Oxígeno",    "O",  "15,999", "16", "2"],
        ["11", "Sodio",      "Na", "22,990", "1",  "3"],
        ["12", "Magnesio",   "Mg", "24,305", "2",  "3"],
        ["13", "Aluminio",   "Al", "26,982", "13", "3"],
        ["14", "Silicio",    "Si", "28,086", "14", "3"],
        ["26", "Hierro",     "Fe", "55,845", "8",  "4"],
        ["29", "Cobre",      "Cu", "63,546", "11", "4"],
        ["47", "Plata",      "Ag", "107,868","11", "5"],
        ["79", "Oro",        "Au", "196,967","11", "6"],
        ["82", "Plomo",      "Pb", "207,200","14", "6"],
        ["92", "Uranio",     "U",  "238,029","--", "7"],
    ],
    "Estados de la materia": [
        ["Estado", "Forma", "Volumen", "Ejemplo"],
        ["Sólido",  "Fija",      "Fijo",     "Hielo, hierro"],
        ["Líquido", "Variable",  "Fijo",     "Agua, mercurio"],
        ["Gas",     "Variable",  "Variable", "Aire, vapor de agua"],
        ["Plasma",  "Variable",  "Variable", "Sol, relámpago"],
    ],
}

NERVIOSO_ODP = [
    ("El Sistema Nervioso",
     "El sistema nervioso es el conjunto de estructuras que reciben, procesan y transmiten "
     "señales entre diferentes partes del cuerpo y entre el cuerpo y el entorno."),
    ("División anatómica",
     "Sistema nervioso central (SNC): encéfalo y médula espinal.\n"
     "Sistema nervioso periférico (SNP): nervios craneales y espinales.\n"
     "El SNP se subdivide en somático (voluntario) y autónomo (involuntario)."),
    ("La neurona",
     "Unidad funcional del sistema nervioso.\n"
     "Partes: cuerpo celular (soma), dendritas y axón.\n"
     "Las neuronas se comunican mediante sinapsis.\n"
     "El cerebro humano contiene ~86 000 millones de neuronas."),
    ("El encéfalo",
     "Cerebro: pensamiento, memoria, lenguaje, movimiento voluntario.\n"
     "Cerebelo: coordinación motora y equilibrio.\n"
     "Tronco del encéfalo: funciones vitales (respiración, ritmo cardíaco).\n"
     "Tálamo: relay sensorial hacia la corteza cerebral."),
    ("Sistema nervioso autónomo",
     "Simpático: respuesta de lucha o huida (aumenta ritmo cardíaco, dilata pupilas).\n"
     "Parasimpático: descanso y digestión (reduce ritmo cardíaco, contrae pupilas).\n"
     "Los dos sistemas actúan de forma complementaria."),
]

ENERGIA_XLSX = {
    "Capacidad instalada (GW)": [
        ["Tecnología", "2000", "2010", "2020", "2023"],
        ["Solar fotovoltaica", "1,4", "40", "714", "1177"],
        ["Eólica terrestre",   "17",  "178", "698", "899"],
        ["Hidroeléctrica",     "748", "990", "1330", "1392"],
        ["Eólica marina",      "0,1", "3,1", "35",  "69"],
        ["Geotérmica",         "8",   "11",  "14",  "15"],
        ["Biomasa",            "27",  "65",  "124", "145"],
        ["Nuclear",            "349", "375", "393", "371"],
    ],
    "Producción eléctrica 2023 (TWh)": [
        ["Fuente", "Producción TWh", "% total mundial"],
        ["Carbón",     "10100", "35,4"],
        ["Gas natural", "6400", "22,5"],
        ["Nuclear",     "2700", "9,5"],
        ["Hidroeléctrica", "4200", "14,7"],
        ["Eólica",      "2300", "8,1"],
        ["Solar",       "1600", "5,6"],
        ["Otras renovables", "600", "2,1"],
        ["Petróleo",    "900",  "3,2"],
    ],
}

HISTORIA_FISICA_RTF = """\
Historia de la Física

La física es la ciencia natural que estudia la materia, la energía y sus interacciones.
Su desarrollo abarca desde la antigüedad hasta la actualidad.

Física clásica (siglos XVII-XIX)

Isaac Newton (1643-1727) formuló las tres leyes del movimiento y la ley de gravitación
universal. Su obra Principia Mathematica (1687) estableció los fundamentos de la mecánica
clásica. La fórmula F = ma describe la relación entre fuerza, masa y aceleración.

James Clerk Maxwell (1831-1879) unificó la electricidad, el magnetismo y la óptica en sus
ecuaciones del electromagnetismo. Predijo la existencia de las ondas electromagnéticas,
incluida la luz visible, y calculó su velocidad: 299 792 km/s.

Ludwig Boltzmann (1844-1906) desarrolló la mecánica estadística y explicó la termodinámica
en términos de comportamiento microscópico de las moléculas.

Revolución cuántica y relatividad (siglo XX)

Max Planck (1858-1947) introdujo en 1900 el concepto de cuanto de energía para explicar
la radiación del cuerpo negro, iniciando la teoría cuántica.

Albert Einstein (1879-1955) publicó en 1905 la teoría especial de la relatividad
(E = mc²) y en 1915 la relatividad general, que revolucionó la comprensión del espacio,
el tiempo y la gravedad. También explicó el efecto fotoeléctrico, base de la mecánica cuántica.

Niels Bohr (1885-1962) propuso el modelo atómico de Bohr (1913), donde los electrones
orbitan el núcleo en niveles de energía cuantizados.

Werner Heisenberg (1901-1976) formuló el principio de incertidumbre: no es posible conocer
simultáneamente con precisión arbitraria la posición y el momento de una partícula.

Física contemporánea

El Modelo Estándar de física de partículas describe con precisión todas las partículas
fundamentales conocidas y tres de las cuatro fuerzas fundamentales: electromagnética,
nuclear fuerte y nuclear débil. El bosón de Higgs fue descubierto en el CERN en 2012.

La gravedad cuántica, que unificaría la relatividad general con la mecánica cuántica,
sigue siendo uno de los grandes problemas sin resolver de la física actual.
"""

COMPUTACION_EPUB = [
    ("Los pioneros", """\
La historia de la computación comienza con Charles Babbage (1791-1871), quien diseñó
la Máquina Analítica, considerada el primer concepto de ordenador de propósito general.
Ada Lovelace escribió el primer algoritmo concebido para ser procesado por esta máquina.

Alan Turing (1912-1954) formalizó el concepto de algoritmo y computación con la Máquina de
Turing (1936). Durante la Segunda Guerra Mundial diseñó la Bombe para descifrar Enigma.
Su artículo de 1950 propuso el Test de Turing como criterio de inteligencia artificial."""),
    ("Los primeros ordenadores", """\
ENIAC (1945) fue el primer computador electrónico de propósito general, pesaba 30 toneladas
y ocupaba 167 m². Realizaba 5000 sumas por segundo usando 17 468 válvulas de vacío.

John von Neumann propuso en 1945 la arquitectura que lleva su nombre: CPU, memoria y
unidad de entrada/salida comunicadas por buses. Esta arquitectura es la base de todos
los ordenadores modernos.

El transistor fue inventado en 1947 en los Bell Labs por Shockley, Bardeen y Brattain.
Sustituyó a las válvulas de vacío, reduciendo tamaño y consumo drásticamente."""),
    ("La era del microprocesador", """\
El circuito integrado fue inventado en 1958 por Jack Kilby (Texas Instruments) e
independientemente por Robert Noyce (Fairchild). Esto permitió miniaturizar la electrónica.

Intel lanzó el 4004 en 1971, el primer microprocesador comercial: 2300 transistores,
740 kHz de frecuencia. En 2023, los procesadores modernos integran más de 100 000 millones
de transistores a 3 nanómetros.

La Ley de Moore (1965) predijo que el número de transistores en un chip se duplicaría
aproximadamente cada dos años. Esta tendencia se mantuvo durante décadas."""),
    ("La era de internet y la IA", """\
Tim Berners-Lee inventó la World Wide Web en 1989 en el CERN. En 1993 Mosaic fue
el primer navegador gráfico popular, democratizando el acceso a internet.

El deep learning revolucionó la inteligencia artificial. En 2012, AlexNet ganó
ImageNet con un 37% de mejora de precisión. En 2022, modelos de lenguaje grandes
como GPT y BERT transformaron el procesamiento del lenguaje natural.

En 2023 los procesadores con IA integrada (NPU) se volvieron estándar en dispositivos
móviles y ordenadores de escritorio, acelerando la inferencia local de modelos de IA."""),
]

ECONOMIA_DOC = """\
Economía: Conceptos Básicos

La economía es la ciencia social que estudia la producción, distribución y consumo de
bienes y servicios. Se divide en microeconomía (comportamiento individual) y macroeconomía
(agregados de la economía nacional e internacional).

El Producto Interior Bruto

El Producto Interior Bruto (PIB) mide el valor de todos los bienes y servicios finales
producidos en un país durante un período, normalmente un año. Es el indicador más utilizado
para medir el tamaño y crecimiento de una economía.
PIB nominal vs PIB real: el PIB real descuenta el efecto de la inflación.
PIB per cápita: PIB dividido entre la población, mide el nivel de vida aproximado.

La Inflación

La inflación es el aumento generalizado y sostenido de los precios de bienes y servicios.
Se mide con el Índice de Precios al Consumo (IPC).
Causas principales: exceso de demanda, aumento de costes de producción, expectativas.
El Banco Central Europeo tiene como objetivo una inflación del 2% anual.

Oferta y Demanda

La ley de la demanda establece que, a mayor precio, menor cantidad demandada (relación
inversa). La ley de la oferta establece que, a mayor precio, mayor cantidad ofertada.
El precio de equilibrio es aquel donde oferta y demanda se igualan.

Los tipos de interés afectan al crédito: tipos bajos estimulan el consumo y la inversión;
tipos altos frenan la inflación pero pueden enfriar la economía.

Sistemas económicos

Economía de mercado: los precios los fija la oferta y la demanda libremente.
Economía planificada: el Estado decide qué, cómo y para quién producir.
Economía mixta: combina mercado libre con intervención estatal. Es el modelo predominante.

El Fondo Monetario Internacional (FMI) y el Banco Mundial son las principales instituciones
económicas internacionales, fundadas en 1944 en Bretton Woods (New Hampshire, EE.UU.).
"""

ASTRONOMIA_PPT = [
    "El Universo: origen y estructura\nEl universo observable tiene un diametro de 93000 millones de anos luz y contiene unos 2 billones de galaxias.",
    "La teoria del Big Bang\nHace 13800 millones de anos el universo comenzo en un estado de densidad y temperatura extremas. La radiacion de fondo de microondas es su eco fosil.",
    "Las galaxias\nLa Via Lactea es una galaxia espiral barrada con unos 200000 millones de estrellas. Se encuentra en el Grupo Local junto a Andromeda.",
    "Las estrellas\nLas estrellas son esferas de plasma que generan energia por fusion nuclear en su nucleo. El Sol es una estrella de tipo G de secuencia principal.",
    "Agujeros negros\nUn agujero negro es una region del espacio donde la gravedad es tan intensa que nada, ni la luz, puede escapar. El horizonte de sucesos marca el punto de no retorno.",
    "Exoplanetas\nMas de 5500 exoplanetas han sido confirmados. La zona habitable es la region donde puede existir agua liquida en la superficie de un planeta.",
]


# ═══════════════════════════════════════════════════════════════════════════
# Helper: OLE2 mínimo (Word .doc y PowerPoint .ppt)
# ═══════════════════════════════════════════════════════════════════════════

def _write_ole2(path: Path, stream_name: str, data: bytes) -> None:
    """
    Crea un contenedor OLE2 (Compound Document Binary) mínimo con un único flujo.
    Compatible con olefile para lectura. Usado para generar .doc y .ppt de prueba.

    Layout de sectores:
      Sector 0  → FAT
      Sector 1  → Directorio (4 entradas × 128 bytes = 512 bytes)
      Sectores 2+ → Datos del flujo
    """
    SECT = 512
    ENDOFCHAIN = 0xFFFFFFFE
    FREESECT   = 0xFFFFFFFF
    FATSECT    = 0xFFFFFFFD

    def pad_to_sector(b: bytes) -> bytes:
        r = len(b) % SECT
        return b + b"\x00" * (SECT - r) if r else b

    # Garantizar que el flujo supera el mini_stream_cutoff (4096 bytes)
    # para que olefile lo busque en los sectores regulares (no en el mini-stream).
    if len(data) <= 4096:
        data = data + b"\x00" * (4097 - len(data))

    data_sz     = len(data)           # tamaño real (para la entrada de directorio)
    data_padded = pad_to_sector(data)
    n_data      = len(data_padded) // SECT

    # FAT: sector 0 = FAT, sector 1 = dir, sectores 2..n+1 = datos
    fat: list[int] = [FATSECT, ENDOFCHAIN]
    for k in range(n_data - 1):
        fat.append(2 + k + 1)
    fat.append(ENDOFCHAIN)
    while len(fat) % 128 != 0:
        fat.append(FREESECT)
    fat_bytes = struct.pack(f"<{len(fat)}I", *fat)

    # Entradas de directorio (128 bytes cada una, formato MS-CFBF §2.6)
    def dir_entry(name: str, obj_type: int, start: int, size: int,
                  child: int = 0xFFFFFFFF, left: int = 0xFFFFFFFF,
                  right: int = 0xFFFFFFFF) -> bytes:
        enc      = name.encode("utf-16-le") if name else b""
        name_len = len(enc) + 2 if enc else 0
        enc      = enc.ljust(64, b"\x00")[:64]
        ts       = b"\x00" * 8
        return (
            enc
            + struct.pack("<H", name_len)          # 0x40: nombre (bytes incl. null)
            + bytes([obj_type, 0])                 # 0x42: tipo, color
            + struct.pack("<III", left, right, child)  # 0x44: siblings + child
            + b"\x00" * 16                         # 0x50: CLSID
            + struct.pack("<I", 0)                 # 0x60: StateBits (4 bytes)
            + ts + ts                              # 0x64: Created + Modified (8+8)
            + struct.pack("<III", start, size, 0)  # 0x74: StartSector + SizeLow + SizeHigh
        )

    root   = dir_entry("Root Entry",  5, ENDOFCHAIN, 0, child=1)
    stream = dir_entry(stream_name,   2, 2,           data_sz)
    dir_sector = root + stream + b"\x00" * 128 + b"\x00" * 128  # 4 entradas × 128

    # Cabecera OLE2 (512 bytes, §2.2)
    header = (
        b"\xD0\xCF\x11\xE0\xA1\xB1\x1A\xE1"  # magic
        + b"\x00" * 16                          # CLSID
        + struct.pack("<H", 0x003E)             # minor version
        + struct.pack("<H", 0x0003)             # major version 3
        + b"\xFE\xFF"                           # byte order (LE)
        + struct.pack("<H", 9)                  # sector pow (2^9 = 512)
        + struct.pack("<H", 6)                  # mini sector pow (2^6 = 64)
        + b"\x00" * 6                           # reservado
        + struct.pack("<I", 0)                  # num dir sectors (0 para v3)
        + struct.pack("<I", 1)                  # num FAT sectors
        + struct.pack("<I", 1)                  # first dir sector = 1
        + struct.pack("<I", 0)                  # transaction signature
        + struct.pack("<I", 0x1000)             # mini stream cutoff = 4096
        + struct.pack("<I", FREESECT)           # first mini FAT sector
        + struct.pack("<I", 0)                  # num mini FAT sectors
        + struct.pack("<I", FREESECT)           # first DIFAT sector
        + struct.pack("<I", 0)                  # num DIFAT sectors
        + struct.pack("<I", 0)                  # DIFAT[0] = sector 0 (FAT)
        + struct.pack("<I", FREESECT) * 108     # DIFAT[1..108] = libre
    )
    assert len(header) == 512, f"Header inesperado: {len(header)} bytes"

    path.write_bytes(header + fat_bytes + dir_sector + data_padded)


# ═══════════════════════════════════════════════════════════════════════════
# Generadores — formatos existentes
# ═══════════════════════════════════════════════════════════════════════════

def write_markdown() -> None:
    path = OUT / "python_lenguaje.md"
    path.write_text(PYTHON_MD, encoding="utf-8")
    print(f"  [OK] {path.name}")


def write_txt_solar() -> None:
    path = OUT / "sistema_solar.txt"
    path.write_text(SOLAR_TXT, encoding="utf-8")
    print(f"  [OK] {path.name}")


def write_pdf_clima() -> None:
    """Genera cambio_climatico.pdf con PyMuPDF."""
    import fitz

    path = OUT / "cambio_climatico.pdf"
    doc = fitz.open()
    page = doc.new_page(width=595, height=842)
    page.insert_textbox(fitz.Rect(72, 72, 523, 770), CLIMA_TXT,
                        fontsize=10.5, fontname="helv", color=(0, 0, 0))
    doc.save(str(path))
    doc.close()
    print(f"  [OK] {path.name}")


def write_docx_ia() -> None:
    """Genera inteligencia_artificial.docx con python-docx."""
    from docx import Document

    path = OUT / "inteligencia_artificial.docx"
    doc = Document()
    for heading, body in IA_CONTENT:
        lvl = 1 if heading == "Inteligencia Artificial" else 2
        doc.add_heading(heading, level=lvl)
        doc.add_paragraph(body)
    doc.save(str(path))
    print(f"  [OK] {path.name}")


def write_pptx_internet() -> None:
    """Genera historia_internet.pptx con python-pptx."""
    from pptx import Presentation

    path = OUT / "historia_internet.pptx"
    prs = Presentation()
    layout = prs.slide_layouts[1]
    for title_text, body_text in INTERNET_PPTX:
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title_text
        slide.placeholders[1].text_frame.text = body_text
    prs.save(str(path))
    print(f"  [OK] {path.name}")


# ═══════════════════════════════════════════════════════════════════════════
# Generadores — formatos nuevos
# ═══════════════════════════════════════════════════════════════════════════

def write_xlsx_energia() -> None:
    """Genera energias_renovables.xlsx con openpyxl."""
    import openpyxl

    path = OUT / "energias_renovables.xlsx"
    wb = openpyxl.Workbook()
    wb.remove(wb.active)  # eliminar hoja vacía por defecto

    for sheet_name, rows in ENERGIA_XLSX.items():
        ws = wb.create_sheet(title=sheet_name[:31])  # max 31 chars
        for row in rows:
            ws.append(row)

    wb.save(str(path))
    print(f"  [OK] {path.name}")


def write_odt_biologia() -> None:
    """Genera biologia_celular.odt con odfpy."""
    from odf.opendocument import OpenDocumentText
    from odf.text import H, P

    path = OUT / "biologia_celular.odt"
    doc = OpenDocumentText()

    for i, (heading, body) in enumerate(BIOLOGIA_ODT):
        level = 1 if i == 0 else 2
        doc.text.addElement(H(outlinelevel=level, text=heading))
        for line in body.split("\n"):
            line = line.strip()
            if line:
                doc.text.addElement(P(text=line))

    doc.save(str(path))
    print(f"  [OK] {path.name}")


def write_ods_tabla() -> None:
    """Genera tabla_periodica.ods con odfpy."""
    from odf.opendocument import OpenDocumentSpreadsheet
    from odf.table import Table, TableCell, TableRow
    from odf.text import P

    path = OUT / "tabla_periodica.ods"
    doc = OpenDocumentSpreadsheet()

    for sheet_name, rows in TABLA_ODS.items():
        table = Table(name=sheet_name)
        for row_data in rows:
            row = TableRow()
            for cell_text in row_data:
                cell = TableCell()
                cell.addElement(P(text=str(cell_text)))
                row.addElement(cell)
            table.addElement(row)
        doc.spreadsheet.addElement(table)

    doc.save(str(path))
    print(f"  [OK] {path.name}")


def write_odp_nervioso() -> None:
    """Genera sistema_nervioso.odp con odfpy."""
    from odf.draw import Frame, Page, TextBox
    from odf.opendocument import OpenDocumentPresentation
    from odf.style import MasterPage, PageLayout, PageLayoutProperties
    from odf.text import P as TextP

    path = OUT / "sistema_nervioso.odp"
    doc = OpenDocumentPresentation()

    # Layout de página requerido
    pl = PageLayout(name="PL1")
    pl.addElement(PageLayoutProperties(
        pagewidth="28cm", pageheight="21cm",
        printorientation="landscape", margin="0cm",
    ))
    doc.automaticstyles.addElement(pl)

    # Master page requerida
    mp = MasterPage(name="MP1", pagelayoutname="PL1")
    doc.masterstyles.addElement(mp)

    for title, content in NERVIOSO_ODP:
        page = Page(masterpagename="MP1")
        doc.presentation.addElement(page)

        # Frame de título
        tf = Frame(width="24cm", height="3cm", x="2cm", y="0.5cm")
        page.addElement(tf)
        tb_title = TextBox()
        tf.addElement(tb_title)
        tb_title.addElement(TextP(text=title))

        # Frame de contenido
        cf = Frame(width="24cm", height="15cm", x="2cm", y="4cm")
        page.addElement(cf)
        tb_content = TextBox()
        cf.addElement(tb_content)
        for line in content.split("\n"):
            line = line.strip()
            if line:
                tb_content.addElement(TextP(text=line))

    doc.save(str(path))
    print(f"  [OK] {path.name}")


def write_rtf_fisica() -> None:
    """Genera historia_fisica.rtf en formato RTF Latin-1."""

    def _rtf_encode(text: str) -> str:
        """Convierte texto UTF-8 a escapes RTF (\\' XX) para caracteres Latin-1 > 127."""
        out: list[str] = []
        for ch in text:
            try:
                byte = ch.encode("latin-1")[0]
            except (UnicodeEncodeError, IndexError):
                byte = ord("?")
            if byte < 0x80:
                out.append(ch)
            else:
                out.append(f"\\'{byte:02x}")
        return "".join(out)

    # Convertir saltos de línea a \par y párrafos a bloques RTF
    lines = HISTORIA_FISICA_RTF.split("\n")
    rtf_lines: list[str] = []
    for line in lines:
        rtf_lines.append(_rtf_encode(line) + r"\par")

    rtf_body = "\n".join(rtf_lines)
    rtf = (
        r"{\rtf1\ansi\ansicpg1252\deff0"
        r"{\fonttbl{\f0\froman\fcharset0 Times New Roman;}}"
        r"\f0\fs22 "
        + rtf_body
        + r"}"
    )

    path = OUT / "historia_fisica.rtf"
    path.write_text(rtf, encoding="latin-1")
    print(f"  [OK] {path.name}")


def write_epub_computacion() -> None:
    """Genera historia_computacion.epub como ZIP manual con XHTML."""

    def _chapter_xhtml(title: str, content: str) -> str:
        paragraphs = "".join(
            f"<p>{line.strip()}</p>"
            for line in content.split("\n")
            if line.strip()
        )
        return (
            "<?xml version='1.0' encoding='utf-8'?>"
            "<!DOCTYPE html>"
            "<html xmlns='http://www.w3.org/1999/xhtml'>"
            f"<head><title>{title}</title></head>"
            f"<body><h1>{title}</h1>{paragraphs}</body>"
            "</html>"
        )

    chapter_ids = [f"ch{i+1}" for i in range(len(COMPUTACION_EPUB))]

    # content.opf
    manifest_items = "\n".join(
        f"<item id='{cid}' href='{cid}.xhtml' media-type='application/xhtml+xml'/>"
        for cid in chapter_ids
    )
    manifest_items += "\n<item id='ncx' href='toc.ncx' media-type='application/x-dtbncx+xml'/>"
    spine_items = "\n".join(
        f"<itemref idref='{cid}'/>" for cid in chapter_ids
    )
    opf = f"""<?xml version='1.0' encoding='utf-8'?>
<package xmlns='http://www.idpf.org/2007/opf' version='2.0' unique-identifier='uid'>
  <metadata xmlns:dc='http://purl.org/dc/elements/1.1/'>
    <dc:title>Historia de la Computacion</dc:title>
    <dc:identifier id='uid'>seekpal-computacion-001</dc:identifier>
    <dc:language>es</dc:language>
  </metadata>
  <manifest>{manifest_items}</manifest>
  <spine toc='ncx'>{spine_items}</spine>
</package>"""

    # toc.ncx
    nav_points = "\n".join(
        f"<navPoint id='np{i+1}' playOrder='{i+1}'>"
        f"<navLabel><text>{title}</text></navLabel>"
        f"<content src='{cid}.xhtml'/>"
        "</navPoint>"
        for i, ((title, _), cid) in enumerate(zip(COMPUTACION_EPUB, chapter_ids))
    )
    ncx = f"""<?xml version='1.0' encoding='utf-8'?>
<ncx xmlns='http://www.daisy.org/z3986/2005/ncx/' version='2005-1'>
  <head><meta name='dtb:uid' content='seekpal-computacion-001'/></head>
  <docTitle><text>Historia de la Computacion</text></docTitle>
  <navMap>{nav_points}</navMap>
</ncx>"""

    path = OUT / "historia_computacion.epub"
    with zipfile.ZipFile(str(path), "w", zipfile.ZIP_DEFLATED) as z:
        z.writestr("mimetype", "application/epub+zip",
                   compress_type=zipfile.ZIP_STORED)
        z.writestr(
            "META-INF/container.xml",
            "<?xml version='1.0'?>"
            "<container version='1.0' xmlns='urn:oasis:names:tc:opendocument:xmlns:container'>"
            "<rootfiles><rootfile full-path='OEBPS/content.opf'"
            " media-type='application/oebps-package+xml'/></rootfiles></container>",
        )
        z.writestr("OEBPS/content.opf", opf)
        z.writestr("OEBPS/toc.ncx", ncx)
        for cid, (title, content) in zip(chapter_ids, COMPUTACION_EPUB):
            z.writestr(f"OEBPS/{cid}.xhtml", _chapter_xhtml(title, content))

    print(f"  [OK] {path.name}")


def write_doc_economia() -> None:
    """
    Genera economia_basica.doc: contenedor OLE2 con flujo WordDocument.

    El flujo comienza con el magic 0xA5EC (FIB de Word) seguido de 510 bytes
    de ceros (FIB mínimo) y luego el texto en UTF-16LE.
    DocExtractor detectará que el piece table falla y usará el fallback heurístico
    (_scan_unicode) que escanea desde el byte 512 buscando cadenas UTF-16LE.
    """
    WORD_MAGIC = 0xA5EC

    text_u16 = ECONOMIA_DOC.encode("utf-16-le")
    fib = struct.pack("<H", WORD_MAGIC) + b"\x00" * 510  # 512 bytes FIB mínimo
    stream_data = fib + text_u16

    _write_ole2(OUT / "economia_basica.doc", "WordDocument", stream_data)
    print("  [OK] economia_basica.doc")


def write_ppt_astronomia() -> None:
    """
    Genera astronomia.ppt: contenedor OLE2 con flujo 'PowerPoint Document'.

    Cada diapositiva se codifica como TextCharsAtom (0x0FA0, UTF-16LE).
    PptExtractor reconoce estos átomos y extrae el texto directamente.
    """
    TEXT_CHARS = 0x0FA0

    stream_data = b""
    for slide_text in ASTRONOMIA_PPT:
        payload = slide_text.encode("utf-16-le")
        stream_data += struct.pack("<HHI", 0, TEXT_CHARS, len(payload)) + payload

    # Padding con átomo desconocido (0x0FFF) para superar el mini_stream_cutoff.
    # _write_ole2 también lo hace automáticamente, pero ser explícito no hace daño.
    _write_ole2(OUT / "astronomia.ppt", "PowerPoint Document", stream_data)
    print("  [OK] astronomia.ppt")


# ═══════════════════════════════════════════════════════════════════════════
# Main
# ═══════════════════════════════════════════════════════════════════════════

if __name__ == "__main__":
    print("Generando corpus de evaluación (12 formatos)...\n")

    # Limpiar corpus anterior
    for f in OUT.glob("*"):
        f.unlink()

    print("-- Formatos de texto / Office moderno ----------------------")
    write_txt_solar()
    write_markdown()
    write_pdf_clima()
    write_docx_ia()
    write_pptx_internet()
    write_xlsx_energia()

    print("\n-- OpenDocument --------------------------------------------")
    write_odt_biologia()
    write_ods_tabla()
    write_odp_nervioso()

    print("\n-- Formatos binarios / especiales --------------------------")
    write_rtf_fisica()
    write_epub_computacion()
    write_doc_economia()
    write_ppt_astronomia()

    print(f"\n[OK] Corpus listo en: {OUT}\n")
    total = 0
    for f in sorted(OUT.iterdir()):
        size = f.stat().st_size
        total += size
        print(f"  {f.name:<45} {size:>9,} bytes")
    print(f"\n  {'TOTAL':<45} {total:>9,} bytes")
