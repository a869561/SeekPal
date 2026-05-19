from __future__ import annotations

from pathlib import Path

import pytest


@pytest.fixture
def tmp_text_file(tmp_path: Path) -> Path:
    p = tmp_path / "sample.txt"
    p.write_text(
        "SeekPal RAG pipeline test.\n"
        "Este fichero contiene texto en español.\n"
        "Probando extracción de texto plano para el año 2025.\n",
        encoding="utf-8",
    )
    return p


@pytest.fixture
def tmp_md_file(tmp_path: Path) -> Path:
    p = tmp_path / "sample.md"
    p.write_text(
        "# Título del documento\n\n"
        "Contenido del año anterior.\n"
        "## Sección 2\n\n"
        "Más contenido SeekPal RAG.\n",
        encoding="utf-8",
    )
    return p


@pytest.fixture
def tmp_pdf_file(tmp_path: Path) -> Path:
    import fitz

    doc = fitz.open()
    try:
        page1 = doc.new_page()
        page1.insert_text((72, 72), "Página 1: SeekPal RAG pipeline.")
        page2 = doc.new_page()
        page2.insert_text((72, 72), "Página 2: contenido adicional para pruebas.")
        out = tmp_path / "sample.pdf"
        doc.save(str(out))
    finally:
        doc.close()
    return out


@pytest.fixture
def tmp_docx_file(tmp_path: Path) -> Path:
    from docx import Document

    doc = Document()
    doc.add_heading("SeekPal Test Document", level=1)
    doc.add_paragraph("Este es el contenido del fichero DOCX de prueba.")
    doc.add_paragraph("Segunda línea para verificar extracción.")
    out = tmp_path / "sample.docx"
    doc.save(str(out))
    return out


@pytest.fixture
def tmp_pptx_file(tmp_path: Path) -> Path:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    layout = prs.slide_layouts[1]

    slide1 = prs.slides.add_slide(layout)
    slide1.shapes.title.text = "SeekPal Slide 1"
    slide1.placeholders[1].text = "Contenido de la diapositiva 1."

    slide2 = prs.slides.add_slide(layout)
    slide2.shapes.title.text = "SeekPal Slide 2"
    slide2.placeholders[1].text = "Contenido de la diapositiva 2."

    out = tmp_path / "sample.pptx"
    prs.save(str(out))
    return out


@pytest.fixture
def tmp_odt_file(tmp_path: Path) -> Path:
    from odf.opendocument import OpenDocumentText
    from odf.text import P

    doc = OpenDocumentText()
    para = P(text="SeekPal RAG ODT test content.")
    doc.text.addElement(para)
    out = tmp_path / "sample.odt"
    doc.save(str(out))
    return out
