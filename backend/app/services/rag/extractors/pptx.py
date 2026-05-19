from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from app.services.rag.extractors.base import BaseExtractor
from app.services.rag.types import ExtractedDoc


class PptxExtractor(BaseExtractor):
    def extract(self, path: Path) -> ExtractedDoc:
        prs = Presentation(str(path))
        parts: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = para.text.strip()
                        if line:
                            parts.append(line)
        text = "\n".join(parts)
        return ExtractedDoc(text=text, page_map=[], extractor="pptx")

    def supported_extensions(self) -> list[str]:
        return [".pptx"]
